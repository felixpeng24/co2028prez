"""
Render dont-die-week-slides.html to a PowerPoint by:
  1. Using headless Chrome to screenshot each slide at a tall viewport (so content fits)
  2. Resizing each PNG down to 1920x1080 (PowerPoint native 16:9)
  3. Building a 16:9 .pptx with one image per slide

Output: dont-die-week-slides.pptx (and a slide-images/ folder with PNGs)
"""

import os, sys, subprocess, time, shutil
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu

HERE = Path(__file__).parent
HTML = HERE / "dont-die-week-slides.html"
IMG_DIR = HERE / "slide-images"
PPTX_OUT = HERE / "dont-die-week-slides.pptx"

NUM_SLIDES = 19
# Chrome --headless=new reserves ~87px of the window height for non-viewport
# overhead, so to get an actual 1920x1080 rendering viewport we need to
# request a window of 1920x1167. The screenshot is then cropped to 1080.
OUT_W, OUT_H = 1920, 1080
CHROME_OVERHEAD = 87
WIN_W, WIN_H = OUT_W, OUT_H + CHROME_OVERHEAD
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

def render_slide(n: int, raw_path: Path, final_path: Path):
    """Render slide N at actual 1920x1080 viewport, save PNG."""
    url = f"file://{HTML}?slide={n}&export=1"
    cmd = [
        CHROME,
        "--headless=new",
        "--disable-gpu",
        "--hide-scrollbars",
        "--force-device-scale-factor=1",
        f"--window-size={WIN_W},{WIN_H}",
        "--virtual-time-budget=8000",
        "--run-all-compositor-stages-before-draw",
        "--default-background-color=00000000",
        f"--screenshot={raw_path}",
        url,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
    if not raw_path.exists() or raw_path.stat().st_size < 1000:
        print(f"  WARN: slide {n} screenshot may have failed", file=sys.stderr)
        if result.stderr:
            print(f"  stderr: {result.stderr[:300]}", file=sys.stderr)
        return False
    img = Image.open(raw_path)
    # Crop to the actual 1920x1080 rendering area (top of the screenshot)
    img = img.crop((0, 0, OUT_W, OUT_H))
    img.save(final_path)
    return True


def main():
    if not HTML.exists():
        sys.exit(f"Missing {HTML}")
    if not Path(CHROME).exists():
        sys.exit(f"Chrome not at {CHROME}")

    IMG_DIR.mkdir(exist_ok=True)
    print(f"Rendering {NUM_SLIDES} slides to {IMG_DIR}/...")

    raw_dir = IMG_DIR / "_raw"
    raw_dir.mkdir(exist_ok=True)
    for n in range(1, NUM_SLIDES + 1):
        raw = raw_dir / f"slide_{n:02d}.png"
        out = IMG_DIR / f"slide_{n:02d}.png"
        if raw.exists(): raw.unlink()
        if out.exists(): out.unlink()
        ok = render_slide(n, raw, out)
        size = out.stat().st_size if out.exists() else 0
        print(f"  slide {n:02d}: {'OK' if ok else 'FAIL'} ({size:,} bytes)")
        time.sleep(0.2)

    print(f"\nBuilding PowerPoint at {PPTX_OUT}...")
    prs = Presentation()
    # 16:9 slide size: 13.333" x 7.5"
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for n in range(1, NUM_SLIDES + 1):
        img_path = IMG_DIR / f"slide_{n:02d}.png"
        if not img_path.exists():
            print(f"  skipping slide {n}, no image")
            continue
        slide = prs.slides.add_slide(blank_layout)
        # Full-bleed image
        slide.shapes.add_picture(
            str(img_path),
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(PPTX_OUT)
    print(f"\nDone. Open {PPTX_OUT}")
    print(f"  Or upload directly to Google Slides via 'File > Import slides'")


if __name__ == "__main__":
    main()
